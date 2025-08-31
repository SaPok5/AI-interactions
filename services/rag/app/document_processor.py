"""Document processing and text extraction"""

import asyncio
import uuid
import hashlib
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
import structlog
import PyPDF2
import docx
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
import markdown
import tiktoken

from .config import settings
from .models import ProcessedDocument, DocumentChunk

logger = structlog.get_logger(__name__)

class DocumentProcessor:
    """Process documents and extract text content"""
    
    def __init__(self):
        self.processing_times = []
        self.supported_formats = settings.supported_formats
        
        # Initialize tokenizer for chunk size calculation
        try:
            self.tokenizer = tiktoken.get_encoding("cl100k_base")
        except:
            self.tokenizer = None
    
    async def process_document(
        self,
        content: bytes,
        filename: str,
        metadata: Dict[str, Any]
    ) -> ProcessedDocument:
        """Process a document and extract text"""
        start_time = datetime.utcnow()
        
        try:
            # Determine file type
            file_extension = filename.split('.')[-1].lower() if '.' in filename else ''
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension}")
            
            # Check file size
            if len(content) > settings.max_file_size_mb * 1024 * 1024:
                raise ValueError(f"File too large: {len(content)} bytes")
            
            # Extract text based on file type
            text_content = await self._extract_text(content, file_extension)
            
            # Generate document ID
            document_id = self._generate_document_id(content, filename)
            
            # Create processed document
            processed_doc = await self.process_text(
                text_content,
                filename,
                metadata,
                document_id
            )
            
            # Record processing time
            processing_time = (datetime.utcnow() - start_time).total_seconds() * 1000
            self.processing_times.append(processing_time)
            if len(self.processing_times) > 1000:
                self.processing_times = self.processing_times[-1000:]
            
            logger.info("Document processed successfully",
                       filename=filename,
                       document_id=document_id,
                       chunks=len(processed_doc.chunks),
                       processing_time_ms=processing_time)
            
            return processed_doc
            
        except Exception as e:
            logger.error("Document processing failed", filename=filename, error=str(e))
            raise
    
    async def process_text(
        self,
        text: str,
        title: str,
        metadata: Dict[str, Any],
        document_id: Optional[str] = None
    ) -> ProcessedDocument:
        """Process raw text content"""
        try:
            if not document_id:
                document_id = self._generate_document_id(text.encode(), title)
            
            # Clean and normalize text
            cleaned_text = self._clean_text(text)
            
            # Split into chunks
            chunks = await self._create_chunks(cleaned_text, document_id, metadata)
            
            # Create processed document
            now = datetime.utcnow()
            processed_doc = ProcessedDocument(
                document_id=document_id,
                title=title,
                content=cleaned_text,
                chunks=chunks,
                metadata=metadata,
                created_at=now,
                updated_at=now
            )
            
            return processed_doc
            
        except Exception as e:
            logger.error("Text processing failed", title=title, error=str(e))
            raise
    
    async def _extract_text(self, content: bytes, file_extension: str) -> str:
        """Extract text from different file formats"""
        try:
            if file_extension == 'pdf':
                return await self._extract_pdf_text(content)
            elif file_extension == 'docx':
                return await self._extract_docx_text(content)
            elif file_extension == 'pptx':
                return await self._extract_pptx_text(content)
            elif file_extension == 'xlsx':
                return await self._extract_xlsx_text(content)
            elif file_extension in ['txt', 'md']:
                return content.decode('utf-8', errors='ignore')
            elif file_extension == 'html':
                return await self._extract_html_text(content)
            else:
                # Try to decode as text
                return content.decode('utf-8', errors='ignore')
                
        except Exception as e:
            logger.error("Text extraction failed", format=file_extension, error=str(e))
            raise
    
    async def _extract_pdf_text(self, content: bytes) -> str:
        """Extract text from PDF"""
        import io
        
        loop = asyncio.get_event_loop()
        
        def extract():
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text_parts = []
            
            for page in pdf_reader.pages:
                text_parts.append(page.extract_text())
            
            return '\n'.join(text_parts)
        
        return await loop.run_in_executor(None, extract)
    
    async def _extract_docx_text(self, content: bytes) -> str:
        """Extract text from DOCX"""
        import io
        
        loop = asyncio.get_event_loop()
        
        def extract():
            doc = docx.Document(io.BytesIO(content))
            text_parts = []
            
            for paragraph in doc.paragraphs:
                text_parts.append(paragraph.text)
            
            return '\n'.join(text_parts)
        
        return await loop.run_in_executor(None, extract)
    
    async def _extract_pptx_text(self, content: bytes) -> str:
        """Extract text from PPTX"""
        import io
        
        loop = asyncio.get_event_loop()
        
        def extract():
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            
            return '\n'.join(text_parts)
        
        return await loop.run_in_executor(None, extract)
    
    async def _extract_xlsx_text(self, content: bytes) -> str:
        """Extract text from XLSX"""
        import io
        
        loop = asyncio.get_event_loop()
        
        def extract():
            workbook = openpyxl.load_workbook(io.BytesIO(content))
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text_parts.append('\t'.join(row_text))
            
            return '\n'.join(text_parts)
        
        return await loop.run_in_executor(None, extract)
    
    async def _extract_html_text(self, content: bytes) -> str:
        """Extract text from HTML"""
        loop = asyncio.get_event_loop()
        
        def extract():
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        
        return await loop.run_in_executor(None, extract)
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        # Remove excessive whitespace
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            if line:
                cleaned_lines.append(line)
        
        # Join lines and normalize spacing
        cleaned_text = '\n'.join(cleaned_lines)
        
        # Remove multiple consecutive newlines
        import re
        cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
        
        return cleaned_text
    
    async def _create_chunks(
        self,
        text: str,
        document_id: str,
        metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """Split text into chunks"""
        chunks = []
        
        # Split text into sentences/paragraphs first
        paragraphs = text.split('\n\n')
        
        current_chunk = ""
        current_start = 0
        chunk_index = 0
        
        for paragraph in paragraphs:
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            
            # Check if adding this paragraph would exceed chunk size
            potential_chunk = current_chunk + '\n\n' + paragraph if current_chunk else paragraph
            
            if self._get_token_count(potential_chunk) <= settings.chunk_size:
                # Add to current chunk
                if current_chunk:
                    current_chunk += '\n\n' + paragraph
                else:
                    current_chunk = paragraph
            else:
                # Save current chunk if it exists
                if current_chunk:
                    chunk = await self._create_chunk(
                        current_chunk,
                        document_id,
                        chunk_index,
                        current_start,
                        current_start + len(current_chunk),
                        metadata
                    )
                    chunks.append(chunk)
                    chunk_index += 1
                
                # Start new chunk
                current_start = text.find(paragraph, current_start)
                current_chunk = paragraph
        
        # Add final chunk
        if current_chunk:
            chunk = await self._create_chunk(
                current_chunk,
                document_id,
                chunk_index,
                current_start,
                current_start + len(current_chunk),
                metadata
            )
            chunks.append(chunk)
        
        return chunks
    
    async def _create_chunk(
        self,
        content: str,
        document_id: str,
        chunk_index: int,
        start_char: int,
        end_char: int,
        metadata: Dict[str, Any]
    ) -> DocumentChunk:
        """Create a document chunk"""
        chunk_id = f"{document_id}_chunk_{chunk_index}"
        
        chunk_metadata = metadata.copy()
        chunk_metadata.update({
            "chunk_index": chunk_index,
            "token_count": self._get_token_count(content),
            "char_count": len(content)
        })
        
        return DocumentChunk(
            chunk_id=chunk_id,
            content=content,
            metadata=chunk_metadata,
            start_char=start_char,
            end_char=end_char
        )
    
    def _get_token_count(self, text: str) -> int:
        """Get token count for text"""
        if self.tokenizer:
            try:
                return len(self.tokenizer.encode(text))
            except:
                pass
        
        # Fallback: rough estimation
        return len(text.split()) * 1.3  # Approximate tokens per word
    
    def _generate_document_id(self, content: bytes, filename: str) -> str:
        """Generate unique document ID"""
        # Create hash from content and filename
        hasher = hashlib.sha256()
        hasher.update(content)
        hasher.update(filename.encode())
        
        return f"doc_{hasher.hexdigest()[:16]}"
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return self.supported_formats
    
    def get_average_processing_time(self) -> float:
        """Get average processing time in milliseconds"""
        if not self.processing_times:
            return 0.0
        return sum(self.processing_times) / len(self.processing_times)
